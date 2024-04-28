from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
import os
import logging

# Ensure logging is configured in your main application entry
logging.basicConfig(level=logging.INFO)

def extract_metadata_and_text(file_path):
    filename = os.path.basename(file_path)
    filetype = os.path.splitext(filename)[1].lower()
    title = author = text = None  # Default to None if not found

    if filetype == '.pdf':
        try:
            with open(file_path, "rb") as f:
                reader = PdfReader(f)
                metadata = reader.metadata
                title = metadata.get('/Title', "Unknown")
                author = metadata.get('/Author', "Unknown")
                text = ' '.join([page.extract_text() for page in reader.pages if page.extract_text()])
        except Exception as e:
            logging.error(f"Error reading PDF file {filename}: {str(e)}")
            text = "Unable to extract text due to an error."

    elif filetype == '.docx':
        try:
            doc = DocxDocument(file_path)
            title = doc.core_properties.title if doc.core_properties.title else "Unknown"
            author = doc.core_properties.author if doc.core_properties.author else "Unknown"
            text = ' '.join([para.text for para in doc.paragraphs if para.text])
        except Exception as e:
            logging.error(f"Error reading DOCX file {filename}: {str(e)}")
            text = "Unable to extract text due to an error."

    elif filetype == '.pptx':
        try:
            pres = Presentation(file_path)
            title = pres.core_properties.title if pres.core_properties.title else "Unknown"
            author = pres.core_properties.author if pres.core_properties.author else "Unknown"
            text = ' '.join([slide.notes_slide.notes_text_frame.text for slide in pres.slides if slide.has_notes_slide and slide.notes_slide.notes_text_frame])
        except Exception as e:
            logging.error(f"Error reading PPTX file {filename}: {str(e)}")
            text = "Unable to extract text due to an error."

    else:
        logging.warning(f"Unsupported file type {filetype} for file {filename}")
        text = "Unsupported file type."

    return {
        "filename": filename,
        "title": title,
        "author": author,
        "filetype": filetype,
        "text": text if text else "No text found or unsupported file type."
    }
