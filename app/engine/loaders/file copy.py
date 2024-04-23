import os
from llama_parse import LlamaParse
from pydantic import BaseModel, validator
from pptx import Presentation
from llama_index.core.readers import SimpleDirectoryReader
import logging
import pdfplumber
from pptx import Presentation
from datetime import datetime

# Enable logging
logging.basicConfig(level=logging.INFO)  # Set to DEBUG for more detailed logs

class DocumentMetadata(BaseModel):
    title: str
    author: str
    creation_date: datetime
    summary: str = None  # Optional field for summary
    content: str

class FileLoaderConfig(BaseModel):
    data_dir: str = "./data"
    use_llama_parse: bool = True

    @validator("data_dir")
    def data_dir_must_exist(cls, v):
        if not os.path.isdir(v):
            raise ValueError(f"Directory '{v}' does not exist")
        return v
    

def llama_parse_parser():
    api_key = os.getenv("LLAMA_CLOUD_API_KEY")
    if api_key is None:
        logging.error("LLAMA_CLOUD_API_KEY is not set. Please configure the environment variable.")
        raise ValueError("LLAMA_CLOUD_API_KEY environment variable is not set.")
    logging.info("Initializing LlamaParse with API key.")
    return LlamaParse(api_key=api_key, result_type="markdown", verbose=True, language="de")

def extract_pdf_metadata_and_text(file_path):
    with pdfplumber.open(file_path) as pdf:
        text = '\n'.join(page.extract_text() for page in pdf.pages if page.extract_text())
        meta = pdf.metadata
        title = meta.get('Title', 'Unknown Title')
        author = meta.get('Author', 'Unknown Author')
        creation_date = meta.get('CreationDate', datetime.now())  # Adjust parsing as necessary
        # Example for adding a summary (you might need a real summarization tool or model here)
        summary = text[:200] + '...'  # Simple truncation for demonstration
        return DocumentMetadata(title=title, author=author, creation_date=creation_date, summary=summary, content=text)

def extract_pptx_metadata_and_text(file_path):
    ppt = Presentation(file_path)
    text = '\n'.join(slide.notes_slide.notes_text_frame.text for slide in ppt.slides if slide.has_notes_slide and slide.notes_slide.notes_text_frame)
    # Simulating metadata extraction
    title = os.path.basename(file_path)
    author = "Unknown Author"
    creation_date = datetime.fromtimestamp(os.path.getctime(file_path))
    summary = text[:200] + '...'  # Simple truncation for demonstration
    return DocumentMetadata(title=title, author=author, creation_date=creation_date, summary=summary, content=text)


def get_file_documents(config: FileLoaderConfig):
    reader = SimpleDirectoryReader(config.data_dir, recursive=True)
    logging.debug(f"Reader type: {type(reader)}")
    if hasattr(reader, 'load_data'):
        logging.debug("load_data is available")
    else:
        logging.debug("load_data is not available")

    if config.use_llama_parse:
        # Assuming LlamaParse is capable of handling these operations or mock it similarly
        reader.file_extractor = {
            ".pdf": extract_pdf_metadata_and_text,
            ".pptx": extract_pptx_metadata_and_text
        }
    else:
        reader.file_extractor = {
            ".pptx": extract_pptx_metadata_and_text  # Only PPTX if LlamaParse is not used
        }

    try:
        documents = []
        data = reader.load_data()
        for file_path, content in data.items():
            doc_metadata = eval(content)  # Assuming content returned as string representation of DocumentMetadata
            documents.append(doc_metadata)
        return documents
    except Exception as e:
        logging.error(f"Failed to load documents: {e}")
        return []

